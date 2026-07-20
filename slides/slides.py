from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# ==========================================
# SLIDE 1: HESA & THE GAP
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout with title
title1 = slide1.shapes.title
title1.text = "HESA & The Policy Gap"

# Diagram: The Pipeline
left = Inches(0.5)
top = Inches(2.0)
width = Inches(2.0)
height = Inches(1.0)

stages = [
    "Stakeholder Briefs\n(544 PDFs)", 
    "Committee Study", 
    "HESA Report\n(Numbered Recs)", 
    "Gov Response\n(120 Days)"
]

for i, stage in enumerate(stages):
    shape = slide1.shapes.add_shape(1, left + Inches(i * 2.3), top, width, height)
    text_frame = shape.text_frame
    text_frame.text = stage
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

# Problem & Promise Text
txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9.0), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.add_paragraph()
p1.text = "The Problem: Thousands of scattered PDFs, unsearchable by theme, perspective, or follow-through."
p1.font.bold = True
p1.font.size = Pt(20)

p2 = tf.add_paragraph()
p2.text = "\nThe Promise:"
p2.font.bold = True
p2.font.size = Pt(20)

p3 = tf.add_paragraph()
p3.text = "• RAG: \"What did nursing organizations ask for on workforce retention?\""
p3.font.size = Pt(18)

p4 = tf.add_paragraph()
p4.text = "• Gap Analysis: \"Which asks made it into recommendations — and which got commitments?\""
p4.font.size = Pt(18)


# ==========================================
# SLIDE 2: DATA ENGINEERING & DISTRIBUTION
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "Content-Based Labeling Reveals True Policy Priorities"
title2.text_frame.paragraphs[0].font.size = Pt(32)

# Pipeline Strip across the top
pipeline_text = "Scrape 562 ➔ Extract/OCR ➔ Clean/Mask ➔ Dedupe ➔ Weak-Label ➔ 250-Word Cards ➔ 217 Gold (Frozen)"
txBox_pipe = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.8))
tf_pipe = txBox_pipe.text_frame
p_pipe = tf_pipe.add_paragraph()
p_pipe.text = pipeline_text
p_pipe.font.size = Pt(16)
p_pipe.font.bold = True
p_pipe.alignment = PP_ALIGN.CENTER

# Topic Distribution Table
rows, cols = 10, 3
table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(2.2), Inches(4.0), Inches(4.0)).table

data = [
    ["Topic", "Count", "%"],
    ["public_health", "176", "34.8%"],
    ["pharmacare", "95", "18.8%"],
    ["womens_health", "80", "15.8%"],
    ["unlabeled", "69", "13.6%"],
    ["workforce", "27", "5.3%"],
    ["mental_health", "24", "4.7%"],
    ["indigenous_health", "14", "2.8%"],
    ["cancer", "11", "2.2%"],
    ["childrens_health", "10", "2.0%"]
]

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table_shape.cell(row_idx, col_idx)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(12)

# Callouts
txBox_callouts = slide2.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(4.5), Inches(3.0))
tf_callouts = txBox_callouts.text_frame
tf_callouts.word_wrap = True

c1 = tf_callouts.add_paragraph()
c1.text = "Data Insight 1:"
c1.font.bold = True
c2 = tf_callouts.add_paragraph()
c2.text = "Labels are two-tier: keyword weak labels for training, and hand-labeled gold (codebook-governed, frozen) for evaluation."

c3 = tf_callouts.add_paragraph()
c3.text = "\nData Insight 2:"
c3.font.bold = True
c4 = tf_callouts.add_paragraph()
c4.text = "1/3 of briefs were COVID-study submissions. Content-based labeling redistributed them across categories because occasion ≠ topic."

prs.save('HESA_Project_Presentation.pptx')
print("Presentation saved as HESA_Project_Presentation.pptx")